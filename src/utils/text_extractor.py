import os
import re
import json
import csv
import email
import logging
import tempfile
import fitz  # PyMuPDF
import google.generativeai as genai
import hashlib
import zipfile
import subprocess
import requests
import numpy as np
from PIL import Image
from pdf2image import convert_from_path
from concurrent.futures import ThreadPoolExecutor, as_completed
import urllib.request
from urllib.parse import urlparse
from typing import Dict, Any, List, Optional, Tuple, Union
from dataclasses import dataclass
import time
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pytesseract
import cv2
from bs4 import BeautifulSoup
import pandas as pd
import xlrd
from tempfile import NamedTemporaryFile
from io import BytesIO
from dotenv import load_dotenv
from .text_cleaner import CleaningOptions, TextCleaner, clean_ocr_text

# === Gemini Setup ===
load_dotenv(dotenv_path=".env.local")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") # Replace with your Gemini API key
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel("gemini-2.5-flash")

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('text_extraction.log')
    ]
)
logger = logging.getLogger(__name__)

@dataclass
class ExtractionResult:
    """Structured result object for text extraction"""
    text: str
    metadata: Dict[str, Any]
    success: bool = True
    error: Optional[str] = None
    processing_time: float = 0.0
    file_hash: Optional[str] = None
    
    def __post_init__(self):
        if not self.success and not self.error:
            self.error = "Unknown error occurred"
    
    def to_tuple(self) -> Tuple[str, Dict[str, Any]]:
        """Convert to backward-compatible tuple format"""
        metadata = self.metadata.copy()
        metadata.update({
            'success': self.success,
            'processing_time': self.processing_time,
            'file_hash': self.file_hash
        })
        if self.error:
            metadata['error'] = self.error
        return (self.text, metadata)

# CleaningOptions imported from text_cleaner.py

# TextCleaner imported from text_cleaner.py

class CacheManager:
    """Simple in-memory cache for extraction results"""
    
    def __init__(self, max_size: int = 100):
        self.cache = {}
        self.max_size = max_size
        self.access_times = {}
    
    def get_hash(self, file_bytes: bytes, filename: str) -> str:
        """Generate hash for caching"""
        content_hash = hashlib.sha256(file_bytes).hexdigest()[:16]
        return f"{filename}_{content_hash}"
    
    def get(self, cache_key: str) -> Optional[ExtractionResult]:
        """Get cached result"""
        if cache_key in self.cache:
            self.access_times[cache_key] = time.time()
            return self.cache[cache_key]
        return None
    
    def set(self, cache_key: str, result: ExtractionResult):
        """Cache result with LRU eviction"""
        if len(self.cache) >= self.max_size:
            # Remove least recently used
            lru_key = min(self.access_times, key=self.access_times.get)
            del self.cache[lru_key]
            del self.access_times[lru_key]
        
        self.cache[cache_key] = result
        self.access_times[cache_key] = time.time()
class PowerPointExtractor:
    def __init__(self, dpi: int = 300, max_workers: int = 4):
        self.dpi = dpi
        self.max_workers = max_workers

    def _extract_text_native(self, pptx_path: str) -> str:
        prs = Presentation(pptx_path)
        slides_text = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if slide_parts:
                slides_text.append(f"Slide {idx}:\n" + "\n".join(slide_parts))
        return "\n\n".join(slides_text)

    def _pptx_to_pdf(self, pptx_path: str, out_dir: str) -> str:
        subprocess.run([
            "soffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_path
        ], check=True)
        base_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
        return os.path.join(out_dir, base_name)

    def _pdf_to_images(self, pdf_path: str, out_dir: str) -> list[str]:
        doc = fitz.open(pdf_path)
        paths = []
        for i, page in enumerate(doc, 1):
            pix = page.get_pixmap(dpi=self.dpi)
            img_path = os.path.join(out_dir, f"slide_{i}.png")
            pix.save(img_path)
            paths.append(img_path)
        return paths

    def _gemini_ocr_image(self, image_path: str, prompt: str = "Extract all text from this slide:") -> str:
        with open(image_path, "rb") as f:
            image_data = f.read()
        response = model.generate_content([
            {"mime_type": "image/png", "data": image_data},
            prompt
        ])
        return clean_ocr_text(response.text)

    def _gemini_pdf_extract(self, pdf_path: str, prompt: str = "Extract slide-wise text from this presentation:") -> str:
        with open(pdf_path, "rb") as f:
            pdf_data = f.read()
        response = model.generate_content([
            {"mime_type": "application/pdf", "data": pdf_data},
            prompt
        ])
        return clean_ocr_text(response.text)

    def _parallel_ocr_images(self, image_paths: list[str]) -> list[str]:
        results = []
        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            futures = [executor.submit(self._gemini_ocr_image, path) for path in image_paths]
            for future in as_completed(futures):
                try:
                    results.append(future.result())
                except Exception as e:
                    results.append(f"[ERROR] {e}")
        return results

    def extract_text_from_pptx_slides(self, pptx_path: str) -> list[str]:
        with tempfile.TemporaryDirectory() as tmp:
            native_text = self._extract_text_native(pptx_path)
            pdf_path = self._pptx_to_pdf(pptx_path, tmp)

            try:
                print("🔍 Using Gemini Vision Pro on full PDF...")
                gemini_text = self._gemini_pdf_extract(pdf_path)
                return [native_text, gemini_text]
            except Exception as e:
                print(f"⚠️ Gemini PDF extract failed: {e}")
                print("🔄 Falling back to image-based OCR...")
                images = self._pdf_to_images(pdf_path, tmp)
                ocr_texts = self._parallel_ocr_images(images)
                return [native_text] + ocr_texts

class PDFExtractor:
    """Optimized PDF extractor with better error handling"""
    
    def __init__(self, cache_manager: Optional[CacheManager] = None):
        self.cache_manager = cache_manager
    
    def _needs_ocr_enhancement(self, text: str, page_area: float = 0) -> bool:
        """Improved OCR need detection"""
        if not text.strip():
            return True
        
        text_length = len(text.strip())
        word_count = len(text.split())
        
        # Multiple heuristics
        indicators = [
            text_length < 50 and page_area > 1000000,  # Little text, large page
            word_count > 0 and text_length / word_count > 50,  # Very long "words"
            len(re.findall(r'[|]{2,}', text)) > 0,  # Multiple vertical bars
            len(re.findall(r'[0O]{3,}', text)) > 0,  # OCR confusion patterns
            text_length > 0 and word_count < text_length * 0.05,  # Very few words
        ]
        
        return sum(indicators) >= 2  # Need at least 2 indicators
    
    def _fuse_hybrid_text(self, digital_text: str, ocr_text: str) -> str:
        """Improved text fusion algorithm"""
        if not digital_text.strip() and not ocr_text.strip():
            return ""
        
        if not digital_text.strip():
            return ocr_text.strip()
        
        if not ocr_text.strip():
            return digital_text.strip()
        
        # Extract important patterns
        email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
        url_pattern = r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
        
        digital_emails = set(re.findall(email_pattern, digital_text, re.IGNORECASE))
        digital_urls = set(re.findall(url_pattern, digital_text))
        ocr_emails = set(re.findall(email_pattern, ocr_text, re.IGNORECASE))
        ocr_urls = set(re.findall(url_pattern, ocr_text))
        
        all_emails = digital_emails.union(ocr_emails)
        all_urls = digital_urls.union(ocr_urls)
        
        # Choose primary text based on quality metrics
        digital_score = self._calculate_text_quality(digital_text)
        ocr_score = self._calculate_text_quality(ocr_text)
        
        if digital_score > ocr_score * 1.2:  # Digital significantly better
            base_text = digital_text
            supplement_text = ocr_text
        elif ocr_score > digital_score * 1.2:  # OCR significantly better
            base_text = ocr_text
            supplement_text = digital_text
        else:  # Similar quality, prefer digital
            base_text = digital_text
            supplement_text = ocr_text
        
        # Add unique valuable content from supplement
        supplement_lines = supplement_text.split('\n')
        unique_content = []
        
        for line in supplement_lines:
            line = line.strip()
            if line and line not in base_text and len(line) > 5:
                # Check if line contains valuable information
                if (any(email in line for email in all_emails) or 
                    any(url in line for url in all_urls) or
                    len(line.split()) > 2):  # Multi-word lines
                    unique_content.append(line)
        
        # Combine texts
        if unique_content:
            fused_text = base_text + "\n\n" + "\n".join(unique_content)
        else:
            fused_text = base_text
        
        # Ensure important data is preserved
        missing_emails = all_emails - set(re.findall(email_pattern, fused_text, re.IGNORECASE))
        missing_urls = all_urls - set(re.findall(url_pattern, fused_text))
        
        if missing_emails or missing_urls:
            appendix = []
            if missing_emails:
                appendix.append("Contact Information: " + ", ".join(missing_emails))
            if missing_urls:
                appendix.append("Web Links: " + ", ".join(missing_urls))
            
            if appendix:
                fused_text += "\n\n" + "\n".join(appendix)
        
        return fused_text.strip()
    
    def _calculate_text_quality(self, text: str) -> float:
        """Calculate text quality score"""
        if not text.strip():
            return 0.0
        
        words = text.split()
        if not words:
            return 0.0
        
        # Various quality metrics
        avg_word_length = sum(len(word) for word in words) / len(words)
        alpha_ratio = sum(1 for char in text if char.isalpha()) / len(text)
        space_ratio = text.count(' ') / len(text)
        sentence_count = len([s for s in re.split(r'[.!?]+', text) if s.strip()])
        
        # Penalize OCR artifacts
        artifact_penalty = (
            text.count('|') * 0.1 +
            len(re.findall(r'[0O]{2,}', text)) * 0.2 +
            len(re.findall(r'[l1]{3,}', text)) * 0.2
        )
        
        quality_score = (
            min(avg_word_length / 5, 1.0) * 0.3 +  # Reasonable word length
            alpha_ratio * 0.3 +  # Alphabetic content
            min(space_ratio * 10, 1.0) * 0.2 +  # Proper spacing
            min(sentence_count / 10, 1.0) * 0.2  # Sentence structure
        ) - artifact_penalty
        
        return max(0.0, quality_score)
    
    def _process_page_optimized(self, doc, page_num: int, enable_ocr: bool, 
                               custom_config: str, combine_digital_and_ocr: bool = True) -> Tuple[int, str, bool]:
        """Enhanced page processing with better error handling"""
        try:
            page = doc.load_page(page_num)
            digital_text = page.get_text() or ""
            ocr_used = False
            ocr_text = ""
            
            # Determine if OCR is needed
            if combine_digital_and_ocr and enable_ocr:
                page_rect = page.rect
                page_area = page_rect.width * page_rect.height
                needs_ocr = self._needs_ocr_enhancement(digital_text, page_area)
                
                if needs_ocr or not digital_text.strip():
                    try:
                        # Enhanced OCR with multiple strategies
                        pix = page.get_pixmap(matrix=fitz.Matrix(2.5, 2.5))  # Higher resolution
                        img_data = pix.tobytes("png")
                        img = Image.open(BytesIO(img_data))
                        
                        # Try multiple OCR approaches
                        ocr_results = self._try_multiple_ocr(img, custom_config)
                        
                        if ocr_results:
                            ocr_text = max(ocr_results, key=len)
                            ocr_used = True
                            logger.debug(f"OCR used for page {page_num + 1}")
                        
                    except Exception as e:
                        logger.warning(f"OCR failed for page {page_num + 1}: {e}")
            
            # Fuse texts intelligently
            if combine_digital_and_ocr and ocr_used and ocr_text.strip():
                final_text = self._fuse_hybrid_text(digital_text, ocr_text)
            else:
                final_text = digital_text
            
            return (page_num, final_text, ocr_used)
            
        except Exception as e:
            logger.error(f"Failed to process page {page_num + 1}: {e}")
            return (page_num, "", False)
    
    def _try_multiple_ocr(self, img: Image.Image, base_config: str) -> List[str]:
        """Try multiple OCR strategies"""
        results = []
        
        # Strategy 1: Basic OCR
        try:
            text = pytesseract.image_to_string(img, config=base_config)
            if text.strip():
                results.append(text)
        except Exception:
            pass
        
        # Strategy 2: Preprocessed OCR
        try:
            img_array = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
            gray = cv2.cvtColor(img_array, cv2.COLOR_BGR2GRAY)
            
            # Multiple preprocessing approaches
            processed_images = [
                cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1],
                cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2),
                cv2.medianBlur(gray, 3)
            ]
            
            for processed in processed_images:
                try:
                    text = pytesseract.image_to_string(processed, config=base_config)
                    if text.strip() and len(text.strip()) > 10:
                        results.append(text)
                        break  # Use first successful result
                except Exception:
                    continue
                    
        except Exception:
            pass
        
        # Strategy 3: Different PSM modes
        if not results or len(results[0].strip()) < 20:
            for psm in [3, 4, 6, 8, 13]:
                try:
                    config = f'--oem 3 --psm {psm}'
                    text = pytesseract.image_to_string(img, config=config)
                    if text.strip() and len(text.strip()) > 10:
                        results.append(text)
                        break
                except Exception:
                    continue
        
        return results
    
    def extract(self, file_bytes: bytes, enable_ocr: bool = True,
                combine_digital_and_ocr: bool = True) -> ExtractionResult:
        """Extract text from PDF with comprehensive error handling"""
        start_time = time.time()
        
        try:
            # Try PyMuPDF first
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            metadata = {
                "pages": doc.page_count,
                "method": "PyMuPDF_hybrid" if combine_digital_and_ocr else "PyMuPDF",
                "ocr_used": False,
                "ocr_pages": [],
                "combine_digital_and_ocr": combine_digital_and_ocr
            }
            
            # Process pages with controlled concurrency
            max_workers = min(8, os.cpu_count() or 1)  # Limit workers
            results = []
            
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = {
                    executor.submit(
                        self._process_page_optimized, 
                        doc, page_num, enable_ocr, "--oem 1 --psm 4", combine_digital_and_ocr
                    ): page_num
                    for page_num in range(doc.page_count)
                }
                
                for future in as_completed(futures, timeout=600):  # 5 minute total timeout
                    try:
                        result = future.result(timeout=120)  # 1 minute per page
                        results.append(result)
                    except Exception as e:
                        page_num = futures[future]
                        logger.warning(f"Page {page_num + 1} processing failed: {e}")
                        results.append((page_num, "", False))
            
            # Sort and assemble results
            results.sort(key=lambda x: x[0])
            page_texts = []
            
            for page_num, page_text, ocr_used in results:
                if page_text.strip():
                    page_texts.append(page_text)
                    if ocr_used:
                        metadata["ocr_pages"].append(page_num)
            
            text = "\n\n".join(page_texts)
            metadata["ocr_used"] = len(metadata["ocr_pages"]) > 0
            doc.close()
            
            processing_time = time.time() - start_time
            
            return ExtractionResult(
                text=text.strip(),
                metadata=metadata,
                success=True,
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"PyMuPDF failed, trying PyPDF2: {e}")
            
            # Fallback to PyPDF2
            try:
                reader = PdfReader(BytesIO(file_bytes))
                page_texts = []
                
                for page in reader.pages:
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        page_texts.append(page_text)
                
                text = "\n\n".join(page_texts)
                processing_time = time.time() - start_time
                
                return ExtractionResult(
                    text=text.strip(),
                    metadata={
                        "pages": len(reader.pages),
                        "method": "PyPDF2_fallback",
                        "ocr_used": False,
                        "fallback_used": True
                    },
                    success=True,
                    processing_time=processing_time
                )
                
            except Exception as e2:
                processing_time = time.time() - start_time
                return ExtractionResult(
                    text="",
                    metadata={"method": "failed"},
                    success=False,
                    error=f"Both PyMuPDF and PyPDF2 failed: {str(e2)}",
                    processing_time=processing_time
                )

class TextExtractor:
    """Main service class with improved architecture"""
    
    def __init__(self, cache_size: int = 100, max_workers: int = None):
        self.cache_manager = CacheManager(cache_size)
        self.max_workers = max_workers or min(8, os.cpu_count() or 1)
        
        # Initialize extractors
        self.pdf_extractor = PDFExtractor(self.cache_manager)
        
        # Register extractors
        self.extractors = {
            ".pdf": self._extract_pdf,
            ".docx": self._extract_docx,
            ".pptx": self._extract_pptx,
            ".png": self._extract_image,
            ".jpg": self._extract_image,
            ".jpeg": self._extract_image,
            ".tiff": self._extract_image,
            ".bmp": self._extract_image,
            ".txt": self._extract_text,
            ".md": self._extract_text,
            ".eml": self._extract_email,
            ".html": self._extract_html,
            ".csv": self._extract_csv,
            ".json": self._extract_json,
            ".xlsx": self._extract_excel,
            ".xls": self._extract_excel,
            ".zip": self._extract_zip,
            ".bin": self._extract_binary,
        }
    
    def extract_text_from_bytes(self, file_bytes: bytes, filename: str, 
                               cleaning_options: Optional[CleaningOptions] = None) -> ExtractionResult:
        """Extract text from file bytes with caching and error handling"""
        start_time = time.time()
        
        # Generate cache key
        cache_key = self.cache_manager.get_hash(file_bytes, filename)
        
        # Check cache
        cached_result = self.cache_manager.get(cache_key)
        if cached_result:
            logger.info(f"Cache hit for {filename}")
            return cached_result
        
        # Set default options
        if cleaning_options is None:
            cleaning_options = CleaningOptions()
        
        try:
            ext = os.path.splitext(filename)[-1].lower()
            
            if ext not in self.extractors:
                return ExtractionResult(
                    text="",
                    metadata={"error": f"Unsupported file type: {ext}"},
                    success=False,
                    error=f"Unsupported file type: {ext}. Supported: {list(self.extractors.keys())}"
                )
            
            # Extract text
            result = self.extractors[ext](file_bytes, cleaning_options)
            
            # Apply cleaning
            if result.success and cleaning_options:
                cleaner = TextCleaner(cleaning_options)
                is_ocr = result.metadata.get("ocr_used", False)
                result.text = cleaner.clean_text(result.text, is_ocr=is_ocr)
            
            # Add processing metadata
            result.processing_time = time.time() - start_time
            result.file_hash = cache_key
            result.metadata.update({
                "filename": filename,
                "file_type": ext,
                "file_size": len(file_bytes),
                "cleaning_applied": cleaning_options is not None
            })
            
            # Cache successful results
            if result.success:
                self.cache_manager.set(cache_key, result)
            
            return result
            
        except Exception as e:
            logger.error(f"Unexpected error extracting text from {filename}: {e}")
            return ExtractionResult(
                text="",
                metadata={"filename": filename},
                success=False,
                error=str(e),
                processing_time=time.time() - start_time
            )
    
    def extract_text_from_url(self, url: str, 
                             cleaning_options: Optional[CleaningOptions] = None) -> ExtractionResult:
        """Download and extract text from URL or local file — skips .bin files"""
        try:
            logger.info(f"Preparing to extract from URL: {url}")
            
            # Handle local file URLs
            if url.startswith('file://'):
                file_path = url[7:]  # Remove 'file://' prefix
                if not os.path.exists(file_path):
                    return ExtractionResult(
                        text="",
                        metadata={"source_url": url},
                        success=False,
                        error=f"Local file not found: {file_path}"
                    )
                
                # Read local file
                with open(file_path, 'rb') as f:
                    file_bytes = f.read()
                
                filename = os.path.basename(file_path)
                logger.info(f"Reading local file: {filename}")
                
            else:
                # Handle remote URLs
                parsed_url = urlparse(url)
                filename = os.path.basename(parsed_url.path) or "document"

                # Early exit if file appears to be a .bin
                if filename.lower().endswith('.bin'):
                    return ExtractionResult(
                        text="",
                        metadata={"source_url": url},
                        success=False,
                        error="Skipping download: '.bin' files are not supported for text extraction."
                    )

                # Download with timeout and user agent
                req = urllib.request.Request(
                    url,
                    headers={'User-Agent': 'Mozilla/5.0 (compatible; TextExtractor/1.0)'}
                )
                
                with urllib.request.urlopen(req, timeout=30) as response:
                    file_bytes = response.read()

                    # Guess extension if missing
                    if '.' not in filename:
                        content_type = response.headers.get('content-type', '').lower()
                        if 'pdf' in content_type:
                            filename += '.pdf'
                        elif 'image' in content_type:
                            filename += '.jpg'
                        else:
                            filename += '.txt'

            # Extract text from bytes
            result = self.extract_text_from_bytes(file_bytes, filename, cleaning_options)
            result.metadata["source_url"] = url

            return result

        except Exception as e:
            logger.error(f"Error downloading/extracting from URL {url}: {e}")
            return ExtractionResult(
                text="",
                metadata={"source_url": url},
                success=False,
                error=str(e)
            )

    def extract_text_from_web_url(self, url: str, 
                                 cleaning_options: Optional[CleaningOptions] = None) -> ExtractionResult:
        """Extract text from web pages with enhanced HTML parsing"""
        try:
            logger.info(f"Extracting text from web page: {url}")
            
            # Download the web page
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            response = requests.get(url, headers=headers, timeout=30)
            response.raise_for_status()
            
            # Parse HTML content
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()
            
            # Extract text with structure preservation
            if cleaning_options and cleaning_options.preserve_structure:
                # Get text with line breaks for structure
                text = soup.get_text(separator='\n', strip=True)
            else:
                # Get plain text
                text = soup.get_text(separator=' ', strip=True)
            
            # Extract metadata
            title = soup.find('title')
            meta_description = soup.find('meta', attrs={'name': 'description'})
            meta_keywords = soup.find('meta', attrs={'name': 'keywords'})
            
            # Apply cleaning if specified
            if cleaning_options:
                cleaner = TextCleaner(cleaning_options)
                text = cleaner.clean_text(text, is_ocr=False)
            
            metadata = {
                "method": "web_parser",
                "title": title.string if title else "",
                "description": meta_description.get('content', '') if meta_description else "",
                "keywords": meta_keywords.get('content', '') if meta_keywords else "",
                "url": url,
                "content_type": response.headers.get('content-type', ''),
                "status_code": response.status_code,
                "content_length": len(response.content)
            }
            
            return ExtractionResult(
                text=text.strip(),
                metadata=metadata,
                success=True,
                processing_time=0.0  # Could add timing if needed
            )
            
        except Exception as e:
            logger.error(f"Error extracting text from web page {url}: {e}")
            return ExtractionResult(
                text="",
                metadata={"url": url, "method": "web_parser"},
                success=False,
                error=str(e)
            )

    def extract_text_from_api_url(self, url: str, 
                                 cleaning_options: Optional[CleaningOptions] = None) -> ExtractionResult:
        """Extract text from API responses (JSON)"""
        try:
            logger.info(f"Extracting text from API: {url}")
            
            # Download the API response
            headers = {
                'User-Agent': 'Mozilla/5.0 (compatible; TextExtractor/1.0)',
                'Accept': 'application/json'
            }
            
            response = requests.get(url, headers=headers, timeout=30)
            response.raise_for_status()
            
            # Parse JSON content
            try:
                data = response.json()
                
                # Convert JSON to readable text
                if cleaning_options and cleaning_options.preserve_formatting:
                    text = json.dumps(data, indent=2, ensure_ascii=False)
                else:
                    text = json.dumps(data, ensure_ascii=False)
                
                # Apply cleaning if specified
                if cleaning_options:
                    cleaner = TextCleaner(cleaning_options)
                    text = cleaner.clean_text(text, is_ocr=False)
                
                metadata = {
                    "method": "api_parser",
                    "url": url,
                    "content_type": response.headers.get('content-type', ''),
                    "status_code": response.status_code,
                    "data_type": type(data).__name__,
                    "data_keys": list(data.keys()) if isinstance(data, dict) else [],
                    "content_length": len(response.content)
                }
                
                return ExtractionResult(
                    text=text.strip(),
                    metadata=metadata,
                    success=True,
                    processing_time=0.0
                )
                
            except json.JSONDecodeError as e:
                # If not valid JSON, treat as plain text
                text = response.text
                
                if cleaning_options:
                    cleaner = TextCleaner(cleaning_options)
                    text = cleaner.clean_text(text, is_ocr=False)
                
                metadata = {
                    "method": "api_parser_text",
                    "url": url,
                    "content_type": response.headers.get('content-type', ''),
                    "status_code": response.status_code,
                    "note": "Response was not valid JSON, treated as text",
                    "content_length": len(response.content)
                }
                
                return ExtractionResult(
                    text=text.strip(),
                    metadata=metadata,
                    success=True,
                    processing_time=0.0
                )
            
        except Exception as e:
            logger.error(f"Error extracting text from API {url}: {e}")
            return ExtractionResult(
                text="",
                metadata={"url": url, "method": "api_parser"},
                success=False,
                error=str(e)
            )
    def extract_text_from_multiple_sources(self, sources: List[Union[str, Tuple[bytes, str]]], 
                                         cleaning_options: Optional[CleaningOptions] = None) -> List[ExtractionResult]:
        """Process multiple sources in parallel"""
        results = []
        
        def process_source(source):
            if isinstance(source, str):  # URL
                return self.extract_text_from_url(source, cleaning_options)
            else:  # (bytes, filename) tuple
                file_bytes, filename = source
                return self.extract_text_from_bytes(file_bytes, filename, cleaning_options)
        
        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            future_to_source = {executor.submit(process_source, source): source for source in sources}
            
            for future in as_completed(future_to_source):
                source = future_to_source[future]
                try:
                    result = future.result()
                    results.append(result)
                    
                    source_id = source if isinstance(source, str) else source[1]
                    status = "✅" if result.success else "❌"
                    logger.info(f"{status} Processed {source_id}")
                    
                except Exception as e:
                    logger.error(f"Unhandled exception for {source}: {e}")
                    results.append(ExtractionResult(
                        text="",
                        metadata={"source": str(source)},
                        success=False,
                        error=str(e)
                    ))
        
        return results
    
    # Extractor methods
    def _extract_pdf(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        return self.pdf_extractor.extract(file_bytes, options.enable_ocr, options.combine_digital_and_ocr)
    
    def _extract_docx(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        try:
            doc = Document(BytesIO(file_bytes))
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            
            # Extract table content
            table_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
            
            text = "\n".join(paragraphs)
            if table_text:
                text += "\n\nTables:\n" + "\n".join(table_text)
            
            return ExtractionResult(
                text=text.strip(),
                metadata={
                    "paragraphs": len(paragraphs),
                    "tables": len(doc.tables),
                    "method": "python-docx"
                },
                success=True
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                metadata={"method": "python-docx"},
                success=False,
                error=str(e)
            )

    def _extract_pptx(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        import tempfile
        import os
        
        # Create a temporary file for the uploaded pptx bytes
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_pptx:
            tmp_pptx.write(file_bytes)
            tmp_pptx_path = tmp_pptx.name

        try:
            # Initialize your extractor
            extractor = PowerPointExtractor(dpi=400, max_workers=6)

            # Extract text from slides
            slide_texts = extractor.extract_text_from_pptx_slides(tmp_pptx_path)

            # Join all slide texts with two newlines between slides
            full_text = "\n\n".join(slide_texts)

            return ExtractionResult(
                text=full_text,
                metadata={"slides": len(slide_texts), "method": "PowerPointExtractor"},
                success=True
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                metadata={"method": "PowerPointExtractor"},
                success=False,
                error=str(e)
            )
        finally:
            # Clean up the temporary pptx file
            if os.path.exists(tmp_pptx_path):
                os.remove(tmp_pptx_path)

    def _extract_image(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        try:
            img = Image.open(BytesIO(file_bytes))
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            results = []
            methods_used = []
            
            # Method 1: Basic Tesseract
            try:
                config = f'--oem 3 --psm 6 -l {options.language}'
                text_basic = pytesseract.image_to_string(img, config=config)
                if text_basic.strip():
                    results.append(text_basic)
                    methods_used.append("Tesseract-basic")
            except Exception as e:
                logger.warning(f"Basic Tesseract failed: {e}")
            
            # Method 2: Preprocessed if basic failed or insufficient
            if not results or len(results[0].strip()) < 50:
                try:
                    img_array = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
                    gray = cv2.cvtColor(img_array, cv2.COLOR_BGR2GRAY)
                    
                    # Multiple preprocessing strategies
                    processed_variants = [
                        cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1],
                        cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2),
                        cv2.medianBlur(gray, 3)
                    ]
                    
                    for i, processed in enumerate(processed_variants):
                        try:
                            text_preproc = pytesseract.image_to_string(processed, config=config)
                            if text_preproc.strip() and len(text_preproc.strip()) > (len(results[0].strip()) if results else 0):
                                results = [text_preproc]  # Replace with better result
                                methods_used = [f"Tesseract-preprocessed-{i}"]
                                break
                        except Exception:
                            continue
                            
                except Exception as e:
                    logger.warning(f"Preprocessed OCR failed: {e}")
            
            # Method 3: Different PSM modes if still insufficient
            if not results or len(results[0].strip()) < 20:
                for psm in [3, 4, 8, 13]:
                    try:
                        psm_config = f'--oem 3 --psm {psm} -l {options.language}'
                        text_psm = pytesseract.image_to_string(img, config=psm_config)
                        if text_psm.strip() and len(text_psm.strip()) > (len(results[0].strip()) if results else 0):
                            results = [text_psm]
                            methods_used = [f"Tesseract-psm{psm}"]
                            break
                    except Exception:
                        continue
            
            # Get confidence scores if available
            confidence_scores = []
            if results:
                try:
                    data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT, config=config)
                    confidence_scores = [int(conf) for conf in data['conf'] if int(conf) > 0]
                except Exception:
                    pass
            
            final_text = results[0] if results else ""
            avg_confidence = sum(confidence_scores) / len(confidence_scores) if confidence_scores else 0
            
            return ExtractionResult(
                text=final_text.strip(),
                metadata={
                    "methods": methods_used,
                    "confidence_scores": confidence_scores,
                    "avg_confidence": avg_confidence,
                    "image_size": img.size,
                    "image_mode": img.mode
                },
                success=bool(final_text.strip())
            )
            
        except Exception as e:
            return ExtractionResult(
                text="",
                metadata={"method": "OCR"},
                success=False,
                error=str(e)
            )
    
    def _extract_text(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            text = ""
            encoding_used = ""
            
            for encoding in encodings:
                try:
                    text = file_bytes.decode(encoding)
                    encoding_used = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if not text:
                # Fallback with error handling
                text = file_bytes.decode('utf-8', errors='ignore')
                encoding_used = 'utf-8-ignore'
            
            return ExtractionResult(
                text=text.strip(),
                metadata={"encoding": encoding_used, "method": "direct"},
                success=True
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                metadata={"method": "direct"},
                success=False,
                error=str(e)
            )
    
    def _extract_email(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        try:
            content_str = file_bytes.decode('utf-8', errors='ignore')
            msg = email.message_from_string(content_str)
            
            text_content = []
            text_content.append(f"Subject: {msg.get('Subject', '')}")
            text_content.append(f"From: {msg.get('From', '')}")
            text_content.append(f"To: {msg.get('To', '')}")
            text_content.append(f"Date: {msg.get('Date', '')}")
            text_content.append("")  # Empty line
            
            body_parts = []
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            body_parts.append(payload.decode('utf-8', errors='ignore'))
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    body_parts.append(payload.decode('utf-8', errors='ignore'))
            
            if body_parts:
                text_content.extend(body_parts)
            
            return ExtractionResult(
                text="\n".join(text_content).strip(),
                metadata={
                    "subject": msg.get('Subject', ''),
                    "from": msg.get('From', ''),
                    "to": msg.get('To', ''),
                    "multipart": msg.is_multipart(),
                    "method": "email"
                },
                success=True
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                metadata={"method": "email"},
                success=False,
                error=str(e)
            )
    
    def _extract_html(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        try:
            content = file_bytes.decode('utf-8', errors='ignore')
            soup = BeautifulSoup(content, "html.parser")
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Extract text with structure preservation
            if options.preserve_structure:
                text = soup.get_text(separator='\n')
            else:
                text = soup.get_text(separator=' ')
            
            # Extract metadata
            title = soup.find('title')
            meta_description = soup.find('meta', attrs={'name': 'description'})
            
            return ExtractionResult(
                text=text.strip(),
                metadata={
                    "title": title.string if title else "",
                    "description": meta_description.get('content', '') if meta_description else "",
                    "parser": "html.parser",
                    "method": "BeautifulSoup"
                },
                success=True
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                metadata={"method": "BeautifulSoup"},
                success=False,
                error=str(e)
            )
    
    def _extract_csv(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        try:
            # Try different encodings and delimiters
            encodings = ['utf-8', 'latin-1', 'cp1252']
            delimiters = [',', ';', '\t', '|']
            
            for encoding in encodings:
                try:
                    content = file_bytes.decode(encoding)
                    # Auto-detect delimiter
                    sample = content[:1024]
                    sniffer = csv.Sniffer()
                    
                    try:
                        delimiter = sniffer.sniff(sample).delimiter
                    except:
                        delimiter = ','  # Default
                    
                    # Read CSV
                    df = pd.read_csv(BytesIO(file_bytes), encoding=encoding, delimiter=delimiter)
                    
                    return ExtractionResult(
                        text=df.to_string(index=False),
                        metadata={
                            "rows": len(df),
                            "columns": len(df.columns),
                            "encoding": encoding,
                            "delimiter": delimiter,
                            "method": "pandas"
                        },
                        success=True
                    )
                except Exception:
                    continue
            
            # Fallback
            raise Exception("Could not parse CSV with any encoding/delimiter combination")
            
        except Exception as e:
            return ExtractionResult(
                text="",
                metadata={"method": "pandas"},
                success=False,
                error=str(e)
            )
    
    def _extract_json(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        try:
            content = file_bytes.decode('utf-8', errors='ignore')
            data = json.loads(content)
            
            if options.preserve_formatting:
                formatted_json = json.dumps(data, indent=2, ensure_ascii=False)
            else:
                formatted_json = json.dumps(data, ensure_ascii=False)
            
            return ExtractionResult(
                text=formatted_json,
                metadata={
                    "valid": True,
                    "type": type(data).__name__,
                    "method": "json"
                },
                success=True
            )
        except json.JSONDecodeError as e:
            return ExtractionResult(
                text="",
                metadata={"method": "json", "valid": False},
                success=False,
                error=f"Invalid JSON: {str(e)}"
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                metadata={"method": "json"},
                success=False,
                error=str(e)
            )
    
    def _extract_excel(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        try:
            with NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            
            try:
                # Try with pandas (modern Excel files)
                try:
                    df_dict = pd.read_excel(tmp_path, sheet_name=None, engine='openpyxl')
                    sheets_text = []
                    
                    for name, df in df_dict.items():
                        sheet_content = f"Sheet: {name}\n{df.to_string(index=False)}"
                        sheets_text.append(sheet_content)
                    
                    return ExtractionResult(
                        text="\n\n".join(sheets_text),
                        metadata={
                            "sheets": len(df_dict),
                            "total_rows": sum(len(df) for df in df_dict.values()),
                            "method": "pandas-openpyxl"
                        },
                        success=True
                    )
                except Exception:
                    # Fallback to xlrd for older files
                    book = xlrd.open_workbook(tmp_path)
                    sheets_text = []
                    
                    for sheet in book.sheets():
                        sheet_data = []
                        for row_idx in range(sheet.nrows):
                            row = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                            sheet_data.append(", ".join(row))
                        
                        sheet_content = f"Sheet: {sheet.name}\n" + "\n".join(sheet_data)
                        sheets_text.append(sheet_content)
                    
                    return ExtractionResult(
                        text="\n\n".join(sheets_text),
                        metadata={
                            "sheets": len(book.sheets()),
                            "method": "xlrd"
                        },
                        success=True
                    )
            finally:
                os.unlink(tmp_path)
                
        except Exception as e:
            return ExtractionResult(
                text="",
                metadata={"method": "excel"},
                success=False,
                error=str(e)
            )

    def _extract_zip(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        """Extract text from ZIP archives by processing all supported files within"""
        try:
            with zipfile.ZipFile(BytesIO(file_bytes), 'r') as zip_file:
                # Get list of files in the archive
                file_list = zip_file.namelist()
                
                # Filter for supported file types
                supported_extensions = {
                    '.pdf', '.docx', '.pptx', '.txt', '.html', '.csv', 
                    '.json', '.xlsx', '.xls', '.eml'
                }
                
                supported_files = [
                    f for f in file_list 
                    if os.path.splitext(f)[1].lower() in supported_extensions
                    and not f.startswith('__MACOSX/')  # Skip macOS metadata
                    and not f.startswith('.')  # Skip hidden files
                ]
                
                if not supported_files:
                    return ExtractionResult(
                        text="",
                        metadata={
                            "method": "zip",
                            "total_files": len(file_list),
                            "supported_files": 0,
                            "file_list": file_list[:10]  # First 10 files for reference
                        },
                        success=False,
                        error="No supported text files found in ZIP archive"
                    )
                
                # Process supported files
                extracted_texts = []
                processed_files = []
                failed_files = []
                
                for filename in supported_files:
                    try:
                        # Read file from ZIP
                        with zip_file.open(filename) as file_in_zip:
                            file_bytes_in_zip = file_in_zip.read()
                        
                        # Extract text using the main service
                        result = self.extract_text_from_bytes(file_bytes_in_zip, filename, options)
                        
                        if result.success and result.text.strip():
                            extracted_texts.append(f"=== {filename} ===\n{result.text}")
                            processed_files.append(filename)
                        else:
                            failed_files.append(filename)
                            
                    except Exception as e:
                        logger.warning(f"Failed to process {filename} in ZIP: {e}")
                        failed_files.append(filename)
                
                if not extracted_texts:
                    return ExtractionResult(
                        text="",
                        metadata={
                            "method": "zip",
                            "total_files": len(file_list),
                            "supported_files": len(supported_files),
                            "processed_files": processed_files,
                            "failed_files": failed_files
                        },
                        success=False,
                        error="Failed to extract text from any files in ZIP archive"
                    )
                
                # Combine all extracted texts
                combined_text = "\n\n".join(extracted_texts)
                
                return ExtractionResult(
                    text=combined_text.strip(),
                    metadata={
                        "method": "zip",
                        "total_files": len(file_list),
                        "supported_files": len(supported_files),
                        "processed_files": processed_files,
                        "failed_files": failed_files,
                        "archive_contents": file_list[:20]  # First 20 files for reference
                    },
                    success=True
                )
                
        except zipfile.BadZipFile:
            return ExtractionResult(
                text="",
                metadata={"method": "zip"},
                success=False,
                error="Invalid or corrupted ZIP file"
            )
        except Exception as e:
            return ExtractionResult(
                text="",
                metadata={"method": "zip"},
                success=False,
                error=f"Error processing ZIP file: {str(e)}"
            )

    def _extract_binary(self, file_bytes: bytes, options: CleaningOptions) -> ExtractionResult:
        """Handle binary files that cannot be parsed as text"""
        try:
            # Try to detect if it's actually a text file with wrong extension
            try:
                # Attempt to decode as text with various encodings
                for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252']:
                    try:
                        text_content = file_bytes.decode(encoding)
                        # If we can decode it and it looks like text, treat it as text
                        if len(text_content) > 0 and any(c.isprintable() for c in text_content[:100]):
                            return ExtractionResult(
                                text=text_content.strip(),
                                metadata={
                                    "method": "binary_as_text",
                                    "encoding": encoding,
                                    "file_size": len(file_bytes),
                                    "note": "Binary file contained readable text content"
                                },
                                success=True
                            )
                    except UnicodeDecodeError:
                        continue
                
                # If we can't decode as text, it's truly binary
                return ExtractionResult(
                    text="",
                    metadata={
                        "method": "binary",
                        "file_size": len(file_bytes),
                        "file_type": "binary_data",
                        "note": "This is a binary file that cannot be parsed as text. Binary files contain non-text data (images, executables, compressed data, etc.) and cannot be processed for text extraction."
                    },
                    success=False,
                    error="Binary file cannot be parsed as text. This file contains non-text data (images, executables, compressed data, etc.) and cannot be processed for text extraction."
                )
                
            except Exception as e:
                return ExtractionResult(
                    text="",
                    metadata={
                        "method": "binary",
                        "file_size": len(file_bytes),
                        "error": str(e)
                    },
                    success=False,
                    error=f"Error analyzing binary file: {str(e)}"
                )
                
        except Exception as e:
            return ExtractionResult(
                text="",
                metadata={"method": "binary"},
                success=False,
                error=f"Unexpected error processing binary file: {str(e)}"
            )


if __name__ == "__main__":
    # CLI test
    import sys
    if len(sys.argv) > 1:
        service = TextExtractor()
        result = service.extract_text_from_url(sys.argv[1], CleaningOptions())
        print(result.text if result.success else f"Error: {result.error}")
